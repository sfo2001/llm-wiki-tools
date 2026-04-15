import pytest
from pathlib import Path


@pytest.fixture
def sample_pdf(tmp_path) -> Path:
    """Minimal valid PDF with text content."""
    try:
        from fpdf import FPDF
        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        pdf.cell(200, 10, txt="Sample Report Title", ln=True)
        pdf.cell(200, 10, txt="This is sample body text for testing.", ln=True)
        path = tmp_path / "sample.pdf"
        pdf.output(str(path))
        return path
    except ImportError:
        pytest.skip("fpdf2 not installed")


@pytest.fixture
def sample_docx(tmp_path) -> Path:
    try:
        from docx import Document
        doc = Document()
        doc.add_heading("Sample Document Title", 0)
        doc.add_paragraph("This is body text.")
        path = tmp_path / "sample.docx"
        doc.save(str(path))
        return path
    except ImportError:
        pytest.skip("python-docx not installed")


@pytest.fixture
def sample_pptx(tmp_path) -> Path:
    try:
        from pptx import Presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = "Slide One Title"
        slide.placeholders[1].text = "Bullet point one"
        path = tmp_path / "sample.pptx"
        prs.save(str(path))
        return path
    except ImportError:
        pytest.skip("python-pptx not installed")


@pytest.fixture
def wiki_dir(tmp_path) -> Path:
    """Minimal wiki with index and two linked pages (no broken links)."""
    d = tmp_path / "wiki"
    d.mkdir()
    (d / "index.md").write_text(
        "# Index\n\n- [[page-a]] — Page A summary\n- [[page-b]] — Page B summary\n"
    )
    (d / "page-a.md").write_text("# Page A\n\nLinks to [[page-b]].\n")
    (d / "page-b.md").write_text("# Page B\n\nNo outbound links.\n")
    return d
