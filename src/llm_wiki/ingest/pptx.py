import logging
import shutil
import subprocess
from pathlib import Path

logger = logging.getLogger(__name__)


def _try_pandoc(path: Path) -> str | None:
    if not shutil.which("pandoc"):
        return None
    result = subprocess.run(
        ["pandoc", "--from", "pptx", "--to", "markdown", str(path)],
        capture_output=True, text=True,
    )
    md = result.stdout.strip()
    return md if result.returncode == 0 and md else None


def _try_python_pptx(path: Path) -> str | None:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        slides = []
        for i, slide in enumerate(prs.slides, start=1):
            title = ""
            bullets: list[str] = []
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                text = shape.text_frame.text.strip()
                if not text:
                    continue
                if shape == slide.shapes.title:
                    title = text
                else:
                    bullets.append(text)
            heading = f"## Slide {i}: {title}" if title else f"## Slide {i}"
            body = "\n".join(f"- {b}" for b in bullets if b)
            slides.append(f"{heading}\n\n{body}" if body else heading)
        return "\n\n".join(slides) if slides else None
    except Exception as e:
        logger.debug("python-pptx failed for %s: %s", path, e)
        return None


def convert_pptx(path: Path) -> tuple[str, str]:
    """Convert PPTX to (backend_name, markdown_body)."""
    md = _try_pandoc(path)
    if md:
        return "pptx.pandoc", md
    md = _try_python_pptx(path)
    if md:
        return "pptx.python-pptx", md
    raise RuntimeError(f"All PPTX backends failed for {path}")
